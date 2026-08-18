#!/usr/bin/env python3
"""Generate tiny, original PDF/image fixtures without third-party copyrighted content."""

from __future__ import annotations

import argparse
import shutil
import subprocess
from pathlib import Path

import pymupdf
from PIL import Image, ImageDraw, ImageFont

ROOT = Path(__file__).resolve().parents[1]
DEFAULT_OUTPUT = ROOT / "tests" / "fixtures"


def _font(size: int = 28) -> ImageFont.FreeTypeFont | ImageFont.ImageFont:
    for candidate in (
        "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
        "/usr/share/fonts/truetype/liberation2/LiberationSans-Regular.ttf",
    ):
        if Path(candidate).is_file():
            return ImageFont.truetype(candidate, size)
    return ImageFont.load_default()


def _make_scan(path: Path) -> None:
    image = Image.new("RGB", (1200, 800), "white")
    draw = ImageDraw.Draw(image)
    font = _font(34)
    small = _font(28)
    draw.text((72, 70), "MOSAICPARSE ORIGINAL SCAN FIXTURE", fill="black", font=font)
    draw.text((72, 150), "Quarter   Revenue   Margin", fill="black", font=small)
    draw.text((72, 205), "Q1 2026   12,345.67  18.25%", fill="black", font=small)
    draw.text((72, 260), "Q2 2026   13,579.20  19.00%", fill="black", font=small)
    draw.rectangle((55, 130, 720, 325), outline="black", width=3)
    draw.line((55, 190, 720, 190), fill="black", width=2)
    draw.line((55, 245, 720, 245), fill="black", width=2)
    image.save(path, format="PNG", optimize=True)


def _make_visuals(output: Path) -> None:
    natural = Image.new("RGB", (960, 640), "#71c7ec")
    draw = ImageDraw.Draw(natural)
    draw.ellipse((690, 55, 820, 185), fill="#f6c85f")
    draw.polygon([(0, 520), (260, 220), (520, 520)], fill="#245274")
    draw.polygon([(280, 520), (650, 165), (960, 520)], fill="#173b5c")
    draw.rectangle((0, 500, 960, 640), fill="#2ec4a6")
    natural.save(output / "natural-scene.png", optimize=True)

    mixed = Image.new("RGB", (960, 640), "#f7f5ef")
    draw = ImageDraw.Draw(mixed)
    draw.rounded_rectangle((50, 45, 910, 595), 22, fill="white", outline="#173b5c", width=4)
    draw.text((90, 80), "MIXED DASHBOARD", fill="#102a43", font=_font(34))
    draw.text((90, 135), "Measured values and visual trend", fill="#102a43", font=_font(24))
    bars = (170, 270, 210, 355, 300)
    for index, height in enumerate(bars):
        left = 110 + index * 130
        draw.rectangle((left, 520 - height, left + 76, 520), fill="#2ec4a6")
    draw.line((90, 520, 850, 520), fill="#102a43", width=4)
    mixed.save(output / "mixed-screenshot.png", optimize=True)
    natural.save(output / "natural-scene.webp", format="WEBP", quality=90)
    natural.save(output / "natural-scene.bmp", format="BMP")
    natural.save(output / "natural-scene.tiff", format="TIFF")


def _make_video(path: Path) -> None:
    ffmpeg = shutil.which("ffmpeg")
    if ffmpeg is None:
        if path.is_file():
            return
        raise RuntimeError("ffmpeg is required to create scene-switch.mp4")
    subprocess.run(
        [
            ffmpeg,
            "-hide_banner",
            "-loglevel",
            "error",
            "-y",
            "-f",
            "lavfi",
            "-i",
            "color=c=#173b5c:s=640x360:d=2:r=12",
            "-f",
            "lavfi",
            "-i",
            "color=c=#2ec4a6:s=640x360:d=2:r=12",
            "-f",
            "lavfi",
            "-i",
            "color=c=#ff786f:s=640x360:d=2:r=12",
            "-filter_complex",
            "[0:v][1:v][2:v]concat=n=3:v=1:a=0,format=yuv420p[v]",
            "-map",
            "[v]",
            "-an",
            "-c:v",
            "libx264",
            "-movflags",
            "+faststart",
            str(path),
        ],
        check=True,
        timeout=60,
    )


def _make_office(output: Path) -> None:
    from docx import Document
    from pptx import Presentation
    from pptx.util import Inches

    image = output / "natural-scene.png"
    document = Document()
    document.add_heading("MosaicParse DOCX fixture", 0)
    document.add_paragraph("This original fixture contains one embedded image.")
    document.add_picture(str(image), width=Inches(4.5))
    document.save(str(output / "embedded-image.docx"))

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "MosaicParse PPTX fixture"
    slide.shapes.add_picture(str(image), Inches(1), Inches(1.7), width=Inches(6))
    presentation.save(str(output / "embedded-image.pptx"))

    video_presentation = Presentation()
    video_slide = video_presentation.slides.add_slide(video_presentation.slide_layouts[5])
    video_slide.shapes.title.text = "Embedded video must be ignored"
    video_slide.shapes.add_movie(
        str(output / "scene-switch.mp4"),
        Inches(1),
        Inches(1.7),
        width=Inches(6),
        height=Inches(3.4),
        poster_frame_image=str(image),
        mime_type="video/mp4",
    )
    video_presentation.save(str(output / "embedded-video.pptx"))


def _new_pdf() -> pymupdf.Document:
    document = pymupdf.open()
    document.set_metadata(
        {
            "title": "MosaicParse original test fixture",
            "author": "MosaicParse contributors",
            "creator": "scripts/generate_fixtures.py",
            "producer": "PyMuPDF",
        }
    )
    return document


def _native_pdf(path: Path) -> None:
    document = _new_pdf()
    page = document.new_page(width=595, height=842)
    page.insert_text((60, 72), "MosaicParse Test Report", fontsize=20)
    page.insert_text((60, 112), "This document has an extractable native text layer.", fontsize=11)
    page.insert_text((60, 150), "Revenue: 12,345.67", fontsize=12)
    page.insert_text((60, 176), "Operating margin: 18.25%", fontsize=12)
    page.insert_text((60, 218), "All words and numbers are original test data.", fontsize=10)
    document.save(path, garbage=4, deflate=True)
    document.close()


def _scan_pdf(path: Path, scan: Path) -> None:
    document = _new_pdf()
    page = document.new_page(width=595, height=397)
    page.insert_image(page.rect, filename=str(scan))
    document.save(path, garbage=4, deflate=True)
    document.close()


def _mixed_pdf(path: Path, scan: Path) -> None:
    document = _new_pdf()
    page = document.new_page(width=595, height=842)
    page.insert_text((60, 72), "Mixed document - native page 1", fontsize=18)
    page.insert_text((60, 112), "Native value: 9,876.54", fontsize=12)
    scanned_page = document.new_page(width=595, height=397)
    scanned_page.insert_image(scanned_page.rect, filename=str(scan))
    document.save(path, garbage=4, deflate=True)
    document.close()


def _columns_pdf(path: Path) -> None:
    document = _new_pdf()
    page = document.new_page(width=595, height=842)
    page.insert_text((60, 60), "Two-column original research fixture", fontsize=17)
    left = pymupdf.Rect(50, 90, 285, 760)
    right = pymupdf.Rect(310, 90, 545, 760)
    page.insert_textbox(
        left,
        "LEFT COLUMN\nSection A\nDemand increased by 7.50%.\nThe reading order begins here.",
        fontsize=11,
    )
    page.insert_textbox(
        right,
        "RIGHT COLUMN\nSection B\nSupply increased by 4.25%.\nThe reading order ends here.",
        fontsize=11,
    )
    document.save(path, garbage=4, deflate=True)
    document.close()


def _table_pdf(path: Path) -> None:
    document = _new_pdf()
    page = document.new_page(width=595, height=842)
    page.insert_text((60, 60), "Original numeric table fixture", fontsize=17)
    x_values = (60, 210, 360, 510)
    y_values = (100, 140, 180, 220)
    for x in x_values:
        page.draw_line((x, y_values[0]), (x, y_values[-1]), width=0.8)
    for y in y_values:
        page.draw_line((x_values[0], y), (x_values[-1], y), width=0.8)
    rows = (
        ("Period", "Revenue", "Margin"),
        ("Q1", "12,345.67", "18.25%"),
        ("Q2", "13,579.20", "19.00%"),
    )
    for row_index, row in enumerate(rows):
        baseline = 126 + row_index * 40
        for column_index, value in enumerate(row):
            page.insert_text((70 + column_index * 150, baseline), value, fontsize=10)
    document.save(path, garbage=4, deflate=True)
    document.close()


def generate(output: Path) -> None:
    output.mkdir(parents=True, exist_ok=True)
    scan = output / "sample-image.png"
    _make_scan(scan)
    _make_visuals(output)
    _make_video(output / "scene-switch.mp4")
    _make_office(output)
    _native_pdf(output / "native-report.pdf")
    _scan_pdf(output / "scanned-report.pdf", scan)
    _mixed_pdf(output / "mixed-report.pdf", scan)
    _columns_pdf(output / "multi-column-research.pdf")
    _table_pdf(output / "table-report.pdf")


def validate(output: Path) -> None:
    expected_pages = {
        "native-report.pdf": 1,
        "scanned-report.pdf": 1,
        "mixed-report.pdf": 2,
        "multi-column-research.pdf": 1,
        "table-report.pdf": 1,
    }
    errors: list[str] = []
    image_path = output / "sample-image.png"
    try:
        with Image.open(image_path) as image:
            image.verify()
    except (FileNotFoundError, OSError) as exc:
        errors.append(f"{image_path.name}: {exc}")
    for name in (
        "natural-scene.png",
        "natural-scene.webp",
        "natural-scene.bmp",
        "natural-scene.tiff",
        "mixed-screenshot.png",
    ):
        try:
            with Image.open(output / name) as image:
                image.verify()
        except (FileNotFoundError, OSError) as exc:
            errors.append(f"{name}: {exc}")
    for name in ("embedded-image.docx", "embedded-image.pptx", "embedded-video.pptx"):
        try:
            import zipfile

            with zipfile.ZipFile(output / name) as archive:
                if "[Content_Types].xml" not in archive.namelist():
                    errors.append(f"{name}: missing OOXML content types")
        except (FileNotFoundError, zipfile.BadZipFile) as exc:
            errors.append(f"{name}: {exc}")
    video = output / "scene-switch.mp4"
    if not video.is_file() or video.stat().st_size < 1_000:
        errors.append("scene-switch.mp4: missing or implausibly small")
    for name, page_count in expected_pages.items():
        path = output / name
        try:
            with pymupdf.open(path) as document:
                if document.page_count != page_count:
                    errors.append(f"{name}: expected {page_count} pages, got {document.page_count}")
        except (FileNotFoundError, RuntimeError) as exc:
            errors.append(f"{name}: {exc}")
    if errors:
        raise SystemExit("fixture validation failed:\n- " + "\n- ".join(errors))


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--output", type=Path, default=DEFAULT_OUTPUT)
    parser.add_argument("--check", action="store_true", help="Validate committed fixtures")
    return parser.parse_args()


def main() -> int:
    args = parse_args()
    if args.check:
        validate(args.output)
    else:
        generate(args.output)
        validate(args.output)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
